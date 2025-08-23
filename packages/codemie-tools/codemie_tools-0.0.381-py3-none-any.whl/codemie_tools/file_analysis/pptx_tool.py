import io
from enum import Enum
from typing import Optional, Type, Any

from langchain_core.language_models import BaseChatModel
from lxml import etree
from pptx import presentation, Presentation
from pptx.shapes.group import GroupShape
from pydantic import BaseModel, Field

from codemie_tools.base.codemie_tool import CodeMieTool
from codemie_tools.file_analysis.tool_vars import PPTX_TOOL

class QueryType(str, Enum):
    TEXT = "Text"
    TEXT_WITH_METADATA = "Text_with_Metadata"
    TOTAL_SLIDES = "Total_Slides"

class PPTXToolInput(BaseModel):
    """
    Defines the schema for the arguments required by PPTXTool.
    """
    slides: list[int] = Field(
        description=(
            "List of slide numbers of a PPTX document to process. "
            "Must be empty to process all slides in a single request. "
            "Slide numbers are 1-based."
        ),
    )
    query: QueryType = Field(
        ...,
        description=(
            "'Text' if the tool must return the Markdown representation of the PPTX slides. "
            "'Text_with_Metadata' if the tool must return the JSON representation of the "
            "PPTX slides with metadata. Preferred if detailed information is needed. "
            "'Total_Slides' if the tool must return the total number of slides in the PPTX "
            "document."
        ),
    )

class PPTXTool(CodeMieTool):
    args_schema: Type[BaseModel] = PPTXToolInput

    name: str = PPTX_TOOL.name
    label: str = PPTX_TOOL.label
    description: str = PPTX_TOOL.description

    # High value to support large PPTX files.
    tokens_size_limit: int = 100_000

    pptx_presentation: Optional[presentation.Presentation] = None

    # This may be used if you want to store or inject a chat model in the future.
    chat_model: Optional[BaseChatModel] = Field(default=None, exclude=True)

    def __init__(self, file_content: bytes, **kwargs: Any) -> None:
        """
        Initialize the PPTXTool with a PPTX as bytes.

        Args:
            file_content (bytes): The raw bytes of the PPTX file.
            **kwargs: Additional keyword arguments to pass along to the super class.
        """
        super().__init__(**kwargs)
        with io.BytesIO(file_content) as file_like:
            self.pptx_presentation = Presentation(file_like)

    def execute(self, slides: list[int], query: QueryType) -> str | dict[str, Any]:
        if not self.pptx_presentation:
            raise ValueError("No PPTX document is loaded. Please provide a valid PPTX.")

        if query == "Total_Slides":
            return str(len(self.pptx_presentation.slides))
        if query.lower().startswith("text"):
            parser = PPTXSimpleParser(self.pptx_presentation)
            dict_output = parser.parse_presentation_to_dict(slides)
            return dict_output if query == "Text_with_Metadata" else parser.convert_dict_to_markdown(dict_output)
        else:
            raise ValueError(
                f"Unknown query '{query}'. Expected one of "
                "['Total_Pages', 'Text', 'Text_with_Metadata']."
            )

class PPTXSimpleParser:
    def __init__(self, pptx_presentation: presentation.Presentation):
        self.pptx_presentation = pptx_presentation

    @staticmethod
    def _parse_text_frame(text_frame):
        """
        Parse the paragraphs/runs in a text frame and return them as a list of dictionaries.
        """
        text_frame_data = []
        for paragraph_index, paragraph in enumerate(text_frame.paragraphs, start=1):
            paragraph_data = {
                "paragraph_index": paragraph_index,
                "runs": []
            }
            for run_index, run in enumerate(paragraph.runs, start=1):
                paragraph_data["runs"].append({
                    "run_index": run_index,
                    "text": run.text
                })
            text_frame_data.append(paragraph_data)
        return text_frame_data

    @staticmethod
    def _parse_table(table):
        """
        Extract information from a table shape and return as a dictionary.
        """
        row_count = len(table.rows)
        col_count = len(table.columns)

        # Gather all cell texts in a 2D array (list of lists)
        table_cells = []
        for r in range(row_count):
            row_data = []
            for c in range(col_count):
                row_data.append(table.cell(r, c).text)
            table_cells.append(row_data)

        return {
            "row_count": row_count,
            "column_count": col_count,
            "cells": table_cells
        }

    @staticmethod
    def _parse_chart(chart):
        """
        Extract basic information from a chart shape (type, title, legend, series).
        Customize as needed to extract more chart details (axes, styling, etc.).
        """
        # Basic chart info
        chart_data = {
            "chart_type": chart.chart_type.name if chart.chart_type else None,
            "has_legend": chart.has_legend
        }

        # Chart title (if present)
        if chart.has_title:
            chart_data["title"] = chart.chart_title.text_frame.text

        # Parse series
        series_list = []
        for s in chart.series:
            series_data = {
                "name": s.name,
                "points": []
            }
            # Each series can have multiple data points
            for p in s.points:
                # category_label is available for category charts, and
                # p.value is the actual numeric value
                point_data = {
                    "category": getattr(p, "category_label", None),
                    "value": getattr(p, "value", None)
                }
                series_data["points"].append(point_data)
            series_list.append(series_data)

        chart_data["series"] = series_list

        return chart_data

    def _parse_shape(self, shape, shape_index):
        """
        Extract detailed information from a single shape and return as a dictionary.
        """
        shape_data = {
            "shape_index": shape_index,
            "name": shape.name,
            "shape_type": shape.shape_type.name if shape.shape_type else "None",
        }
        # Text frames
        if shape.has_text_frame:
            shape_data["text_frame"] = self._parse_text_frame(shape.text_frame)

        # # Pictures are not supported in this version

        # Tables
        if shape.has_table:
            shape_data["table"] = self._parse_table(shape.table)

        # Charts
        if shape.has_chart:
            shape_data["chart"] = self._parse_chart(shape.chart)

        if isinstance(shape, GroupShape):
            group_slides = []
            for shape_index, shape in enumerate(shape.shapes, start=1):
                content = self._parse_shape(shape, shape_index)
                group_slides.append(content)
            shape_data["group_content"] = {"shapes": group_slides}

        return shape_data

    def _parse_modern_comments(self, slide):
        """
        Parse 'modern' comments from the slide's underlying XML (p188:cmLst).
        Returns a list of dictionaries with comment data, including replies and reactions.
        """

        # 1. Locate the XML part containing modern comments
        modern_comment_part = self._find_modern_comment_part(slide)
        if not modern_comment_part:
            return []

        # 2. Parse the XML root
        comment_root = etree.fromstring(modern_comment_part.blob)
        nsmap = {
            "p188": "http://schemas.microsoft.com/office/powerpoint/2018/8/main",
            "a":    "http://schemas.openxmlformats.org/drawingml/2006/main",
            "p223": "http://schemas.microsoft.com/office/powerpoint/2022/03/main",
        }

        # 3. Extract all top-level <p188:cm> nodes
        cm_nodes = comment_root.findall("p188:cm", namespaces=nsmap)
        comments_data = []
        for cm_el in cm_nodes:
            comment_data = {
                "created": cm_el.get("created", ""),  # e.g. "2025-01-17T20:48:57.138"
                "text": self._extract_text(cm_el.find("p188:txBody", namespaces=nsmap), nsmap),
                "reactions": self._parse_reactions(cm_el.find("p188:extLst", namespaces=nsmap), nsmap),
                "replies": self._parse_replies(cm_el, nsmap),
            }
            comments_data.append(comment_data)

        return comments_data

    @staticmethod
    def _find_modern_comment_part(slide):
        """
        Search the slide relationships for a 'modern comments' part.
        Typical rel_type might contain 'comments' and '2018'.
        Returns the part if found; otherwise None.
        """
        for rel in slide.part.rels:
            rel_type = slide.part.rels[rel].reltype
            if "comments" in rel_type and "2018" in rel_type:
                return slide.part.rels[rel].target_part
        return None

    @staticmethod
    def _extract_text(tx_body_el, nsmap):
        """
        Given the <p188:txBody> element, collect all <a:t> text within it
        and return as a single string (lines joined by newline).
        """
        if tx_body_el is None:
            return ""
        text_nodes = tx_body_el.findall(".//a:t", namespaces=nsmap)
        return "\n".join(n.text.strip() for n in text_nodes if n.text).strip()


    def _parse_replies(self, cm_el, nsmap):
        """
        Given a top-level <p188:cm> element, parse any <p188:reply> children
        and return a list of reply dictionaries (created, text, reactions).
        """
        replies_data = []
        reply_list_el = cm_el.find("p188:replyLst", namespaces=nsmap)
        if reply_list_el is None:
            return replies_data

        for reply_el in reply_list_el.findall("p188:reply", namespaces=nsmap):
            replies_data.append({
                "created": reply_el.get("created", ""),
                "text": self._extract_text(reply_el.find("p188:txBody", namespaces=nsmap), nsmap),
                "reactions": self._parse_reactions(reply_el.find("p188:extLst", namespaces=nsmap), nsmap),
            })
        return replies_data


    @staticmethod
    def _parse_reactions(ext_lst_el, nsmap):
        """
        Given an <extLst> element (p188:extLst), parse <p223:rxn> nodes
        to extract reaction 'type' (e.g. 👍).
        Ignores reaction timestamps/author IDs for brevity.
        Returns a list of reaction dictionaries.
        """
        reactions = []
        if ext_lst_el is None:
            return reactions

        rxn_nodes = ext_lst_el.findall(".//p223:rxn", namespaces=nsmap)
        for rxn_el in rxn_nodes:
            rxn_type = rxn_el.get("type")  # e.g. "👍"
            if not rxn_type:
                continue
            # Optionally handle multiple <p223:instance> children if needed
            reactions.append({"type": rxn_type})
        return reactions

    def _parse_slide(self, slide, slide_index):
        """
        Parse all shapes on a slide and return them in a dictionary structure.
        Also checks if there are any notes for the slide.
        """
        slide_data = {
            "slide_index": slide_index,
            "shapes": []
        }

        # Parse shapes on the slide
        for shape_index, shape in enumerate(slide.shapes, start=1):
            shape_data = self._parse_shape(shape, shape_index)
            slide_data["shapes"].append(shape_data)

        # Parse notes, if present
        if slide.notes_slide is not None:
            # notes_slide can contain multiple shapes, but typically
            # the main notes text is in notes_slide.notes_text_frame
            notes_frame = slide.notes_slide.notes_text_frame
            if notes_frame:
                slide_data["notes"] = self._parse_text_frame(notes_frame)
            else:
                slide_data["notes"] = []
        else:
            slide_data["notes"] = []

        slide_data["comments"] = self._parse_modern_comments(slide)

        return slide_data

    def parse_presentation_to_dict(self, slides: list[int]) -> dict:
        """
        Load the PPTX presentation from a byte array, parse only the slides
        whose indexes are in the 'slides' list, and return a dictionary representation.

        :param slides: A list of 1-based slide indexes to parse.
        :return: Dictionary with "slides" key, each item detailing the parsed slides.
        """
        presentation_data = {"slides": []}
        presentation_slides = self.pptx_presentation.slides

        for slide_index, slide in enumerate(presentation_slides, start=1):
            if (not slides) or slide_index in slides:
                slide_data = self._parse_slide(slide, slide_index)
                presentation_data["slides"].append(slide_data)

        return presentation_data

    @classmethod
    def convert_dict_to_markdown(cls, presentation_data: dict) -> str:
        """
        Converts the dictionary representation of a PPTX to Markdown format.

        Args:
            presentation_data (dict): Dictionary containing the parsed PPTX data.
            cls: The class itself.

        Returns:
            str: Markdown-formatted string.
        """
        md_lines = []

        for slide in presentation_data.get("slides", []):
            md_lines.append(_format_slide_header(slide))

            _collect_shapes(md_lines, slide)

            if slide.get("notes"):
                md_lines.append(_format_notes(slide["notes"]))

            if slide.get("comments"):
                md_lines.append(_format_comments(slide["comments"]))

            md_lines.append("\n---")

        return "\n".join(md_lines)

def _collect_shapes(md_lines, shapes, nesting_level = 0):
    for shape in shapes.get("shapes", []):
        md_lines.append(format_shape_header(shape, nesting_level))

        if "text_frame" in shape:
            md_lines.append(_format_text_frame(shape["text_frame"], nesting_level))

        if "table" in shape:
            md_lines.append(_format_table(shape["table"]))

        if "chart" in shape:
            md_lines.append(_format_chart(shape["chart"]))

        if "group_content" in shape:
            _collect_shapes(md_lines, shape["group_content"], nesting_level + 1)

def _format_slide_header(slide):
    return f"# Slide {slide['slide_index']}"

def format_shape_header(shape, nesting_level = 0):
    return f"##{"#" * nesting_level} Shape {shape['shape_index']} - {shape['name']} ({shape['shape_type']})"

def _format_text_frame(text_frame, nesting_level = 0):
    text_lines = ["".join(run["text"] for run in paragraph["runs"]) for paragraph in text_frame]
    return f"###{"#" * nesting_level} Text\n{"\n\n".join(text_lines) if text_lines else "\n"}\n"

def _format_table(table, nesting_level = 0):
    if not table["cells"]:
        return ""

    table_lines = [f"###{"#" * nesting_level} Table"]

    # Add table header
    header_row = table["cells"][0]
    table_lines.append(f"| {' | '.join(header_row)} |")
    table_lines.append(f"|{'|'.join(['---'] * len(header_row))}|")

    # Add table rows
    for row in table["cells"][1:]:
        table_lines.append(f"| {' | '.join(row)} |")

    return "\n".join(table_lines)

def _format_chart(chart, nesting_level = 0):
    chart_lines = [f"###{"#" * nesting_level} Chart\n",
                   f"- Chart Type: {chart['chart_type']}\n",
                   f"- Title: {chart.get('title', 'None')}\n"]

    for series in chart.get("series", []):
        chart_lines.append(f"  - Series: {series['name']}")
        for point in series.get("points", []):
            chart_lines.append(f"    - Category: {point['category']}, Value: {point['value']}")

    return "\n".join(chart_lines)

def _format_notes(notes):
    notes_lines = ["## Notes"] + [
        f"- {run['text']}" for paragraph in notes for run in paragraph["runs"]
    ]
    return "\n".join(notes_lines)

def _format_comments(comments):
    comments_lines = ["## Comments"]
    for comment in comments:
        comments_lines.append(f"- Created: {comment['created']}")
        comments_lines.append(f"  Text: {comment['text']}")
        comments_lines.extend(_collect_comment_reactions(comment))
        comments_lines.extend(_collect_comment_replies(comment))
    return "\n".join(comments_lines)

def _collect_comment_reactions(comment):
    comments_lines = []
    if comment.get("reactions"):
        comments_lines.append("  Reactions:")
        for reaction in comment["reactions"]:
            comments_lines.append(f"    - {reaction['type']}")
    return comments_lines

def _collect_comment_replies(comment):
    comments_lines = []
    if comment.get("replies"):
        comments_lines.append("  Replies:")
        for reply in comment["replies"]:
            comments_lines.append(f"    - Created: {reply['created']}")
            comments_lines.append(f"      Text: {reply['text']}")
            if reply.get("reactions"):
                comments_lines.append("      Reactions:")
                for reaction in reply["reactions"]:
                    comments_lines.append(f"        - {reaction['type']}")
    return comments_lines
