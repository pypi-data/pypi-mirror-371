"""
Document Processor
"""

import os
import io
import zipfile
import re
from typing import List, Dict, Any, Optional, Tuple, Callable
from PIL import Image
import traceback
import logging
import base64
from pathlib import Path

try:
    import pdfplumber
    from pdf2image import convert_from_path
    import pytesseract
    PDF_LIBS_AVAILABLE = True
except ImportError:
    PDF_LIBS_AVAILABLE = False
    pdfplumber, convert_from_path, pytesseract = None, None, None

import docx
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import pandas as pd
import openpyxl

from fairlib.core.types import Document

logger = logging.getLogger(__name__)

try:
    from fairlib import settings
    CONFIG_AVAILABLE = True
except ImportError:
    CONFIG_AVAILABLE = False
    logger.warning("Could not import settings from fairlib!")


class DocumentProcessor:
    """
    Handles document processing, text extraction, and chunking using FAIR-LLM configuration
    """
    
    def __init__(self, config: Optional[Dict[str, Any]] = None):
        """
        Initialize the document processor
        
        Args:
            config: Optional configuration dictionary
        """
        self.config = config or {}
        
        # Load configuration from settings if available
        if CONFIG_AVAILABLE and hasattr(settings, 'rag_system'):
            rag_config = settings.rag_system
            
            # Configuration
            self.files_directory = self.config.get(
                "files_directory",
                rag_config.paths.files_directory
            )
            self.supported_extensions = self.config.get(
                "supported_extensions",
                set(rag_config.document_processing.supported_extensions)
            )
            self.max_chunk_chars = self.config.get(
                "max_chunk_chars",
                rag_config.document_processing.max_chunk_chars
            )
            
            # OCR configuration
            self.enable_ocr = self.config.get(
                "enable_ocr",
                rag_config.document_processing.enable_ocr and PDF_LIBS_AVAILABLE
            )
            self.ocr_dpi = self.config.get(
                "ocr_dpi",
                rag_config.document_processing.ocr_dpi
            )
            self.ocr_grayscale = self.config.get("ocr_grayscale", True)
        
        # Ensure supported_extensions is a set
        if isinstance(self.supported_extensions, list):
            self.supported_extensions = set(self.supported_extensions)
        
        # Processing statistics
        self.stats = {
            "files_processed": 0,
            "extraction_errors": 0,
            "total_chunks": 0,
            "ocr_attempts": 0,
            "ocr_successes": 0
        }
        
        # Ensure files directory exists
        self.ensure_files_directory()
        
        logger.info(f"DocumentProcessor initialized with directory: {self.files_directory}")
    
    def ensure_files_directory(self) -> None:
        """Ensure the files directory exists"""
        os.makedirs(self.files_directory, exist_ok=True)
        logger.info(f"Ensured files directory exists: {self.files_directory}")
    
    def save_uploaded_file(self, file_data: Dict[str, Any]) -> str:
        """
        Save an uploaded file from pipeline data
        
        Args:
            file_data: Dictionary containing file information
                      Expected keys: 'name', 'content' or 'data'
        
        Returns:
            Path to the saved file
        """
        filename = file_data.get('name', 'unnamed_file')
        file_path = os.path.join(self.files_directory, filename)
        
        # Get file content
        content = file_data.get('content') or file_data.get('data')
        
        if isinstance(content, str):
            # If content is base64 encoded
            try:
                content = base64.b64decode(content)
            except Exception as e:
                logger.error(f"Failed to decode base64 content: {e}")
                raise ValueError(f"Invalid base64 content: {e}")
        elif isinstance(content, bytes):
            # Already bytes
            pass
        else:
            raise ValueError(f"Unsupported file content type: {type(content)}")
        
        # Create subdirectories if needed
        os.makedirs(os.path.dirname(file_path), exist_ok=True)
        
        with open(file_path, "wb") as f:
            f.write(content)
        
        logger.info(f"Saved uploaded file to {file_path}")
        return file_path
    
    def extract_text_with_fallback(self, file_path: str) -> List[Dict[str, Any]]:
        """
        Extract text from various file formats
        
        Args:
            file_path: Path to the file to process
            
        Returns:
            List of segments with text and metadata
        """
        logger.info(f"Extracting text from: {file_path}")
        
        filename = os.path.basename(file_path)
        ext = os.path.splitext(file_path)[1].lower()
        
        if ext not in self.supported_extensions:
            logger.warning(f"Unsupported file extension: {ext}")
            return []
        
        # Route to appropriate extraction method
        extraction_methods = {
            ".pdf": self._extract_pdf,
            ".docx": self._extract_docx,
            ".pptx": self._extract_pptx,
            ".xlsx": self._extract_xlsx,
            ".txt": self._extract_text,
            ".csv": self._extract_text,
            ".sql": self._extract_text
        }
        
        method = extraction_methods.get(ext, self._extract_text)
        
        try:
            segments = method(file_path, filename)
            self.stats["files_processed"] += 1
            return segments
        except Exception as e:
            logger.error(f"Error extracting text from {filename}: {traceback.format_exc()}")
            self.stats["extraction_errors"] += 1
            return []
    
    def _extract_pdf(self, file_path: str, filename: str) -> List[Dict[str, Any]]:
        """Extract text from PDF files"""
        if not PDF_LIBS_AVAILABLE:
            logger.warning("PDF processing libraries not available. Skipping PDF.")
            return []
        
        segments = []
        
        try:
            with pdfplumber.open(file_path) as pdf:
                for page_index, page in enumerate(pdf.pages):
                    text = page.extract_text(x_tolerance=1, y_tolerance=3) or ""
                    
                    # Try OCR if no text found and OCR is enabled
                    if not text.strip() and self.enable_ocr and pytesseract:
                        text = self._ocr_pdf_page(file_path, page_index)
                    
                    if text.strip():
                        segments.append({
                            "text": text,
                            "metadata": {
                                "filename": filename,
                                "page": page_index + 1,
                                "type": "pdf"
                            }
                        })
        except Exception as e:
            logger.error(f"PDF processing error for {filename}: {e}")
            raise
        
        return segments
    
    def _ocr_pdf_page(self, file_path: str, page_index: int) -> str:
        """Perform OCR on a PDF page"""
        self.stats["ocr_attempts"] += 1
        try:
            imgs = convert_from_path(
                file_path,
                dpi=self.ocr_dpi,
                first_page=page_index + 1,
                last_page=page_index + 1,
                grayscale=self.ocr_grayscale
            )
            if imgs:
                text = pytesseract.image_to_string(imgs[0])
                self.stats["ocr_successes"] += 1
                return text
        except Exception as ocr_e:
            logger.warning(f"OCR failed for page {page_index + 1}: {ocr_e}")
        return ""
    
    def _extract_docx(self, file_path: str, filename: str) -> List[Dict[str, Any]]:
        """Extract text from DOCX files"""
        segments = []
        
        try:
            doc = docx.Document(file_path)
            full_text_parts = [par.text for par in doc.paragraphs if par.text.strip()]
            
            # Extract text from images if OCR is enabled
            if self.enable_ocr and pytesseract:
                ocr_texts = self._extract_docx_images(file_path)
                if ocr_texts:
                    full_text_parts.append("\n--- OCR Text from Images ---\n" + "\n".join(ocr_texts))
            
            if full_text_parts:
                segments.append({
                    "text": "\n".join(full_text_parts),
                    "metadata": {
                        "filename": filename,
                        "type": "docx"
                    }
                })
        except Exception as e:
            logger.error(f"DOCX processing error for {filename}: {e}")
            raise
        
        return segments
    
    def _extract_docx_images(self, file_path: str) -> List[str]:
        """Extract text from images in DOCX files using OCR"""
        ocr_texts = []
        
        try:
            with zipfile.ZipFile(file_path) as z_docx:
                for name in z_docx.namelist():
                    if name.startswith("word/media/"):
                        self.stats["ocr_attempts"] += 1
                        try:
                            img_bytes = z_docx.read(name)
                            img = Image.open(io.BytesIO(img_bytes))
                            text = pytesseract.image_to_string(img)
                            if text.strip():
                                ocr_texts.append(text)
                                self.stats["ocr_successes"] += 1
                        except Exception as img_e:
                            logger.warning(f"Could not OCR image {name}: {img_e}")
        except Exception as zip_e:
            logger.warning(f"Could not process zip for docx images: {zip_e}")
        
        return ocr_texts
    
    def _extract_pptx(self, file_path: str, filename: str) -> List[Dict[str, Any]]:
        """Extract text from PPTX files"""
        segments = []
        
        try:
            prs = Presentation(file_path)
            for slide_index, slide in enumerate(prs.slides):
                slide_text_parts = []
                
                # Extract text from shapes
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        slide_text_parts.append(shape.text.strip())
                
                # Extract text from images if OCR is enabled
                if self.enable_ocr and pytesseract:
                    ocr_texts = self._extract_pptx_slide_images(slide)
                    if ocr_texts:
                        slide_text_parts.append("\n--- OCR Text from Images ---\n" + "\n".join(ocr_texts))
                
                if slide_text_parts:
                    segments.append({
                        "text": "\n".join(filter(None, slide_text_parts)),
                        "metadata": {
                            "filename": filename,
                            "slide": slide_index + 1,
                            "type": "pptx"
                        }
                    })
        except Exception as e:
            logger.error(f"PPTX processing error for {filename}: {e}")
            raise
        
        return segments
    
    def _extract_pptx_slide_images(self, slide) -> List[str]:
        """Extract text from images in a PowerPoint slide"""
        ocr_texts = []
        
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                self.stats["ocr_attempts"] += 1
                try:
                    image_bytes = shape.image.blob
                    img = Image.open(io.BytesIO(image_bytes))
                    text = pytesseract.image_to_string(img)
                    if text.strip():
                        ocr_texts.append(text)
                        self.stats["ocr_successes"] += 1
                except Exception as img_e:
                    logger.warning(f"Could not OCR image: {img_e}")
        
        return ocr_texts
    
    def _extract_xlsx(self, file_path: str, filename: str) -> List[Dict[str, Any]]:
        """Extract text from XLSX files"""
        segments = []
        
        try:
            xls = pd.ExcelFile(file_path)
            for sheet_name in xls.sheet_names:
                try:
                    df = pd.read_excel(xls, sheet_name=sheet_name, engine='openpyxl')
                    if not df.empty:
                        sheet_text = df.to_string(index=False, na_rep='NaN')
                    else:
                        sheet_text = f"Sheet '{sheet_name}' is empty."
                    
                    segments.append({
                        "text": sheet_text,
                        "metadata": {
                            "filename": filename,
                            "sheet": sheet_name,
                            "type": "xlsx"
                        }
                    })
                except Exception as sheet_e:
                    logger.warning(f"Error reading sheet '{sheet_name}' from {filename}: {sheet_e}")
        except Exception as e:
            logger.error(f"XLSX processing error for {filename}: {e}")
            raise
        
        return segments
    
    def _extract_text(self, file_path: str, filename: str) -> List[Dict[str, Any]]:
        """Extract text from plain text files (txt, csv, sql, etc.)"""
        segments = []
        
        try:
            # Try UTF-8 first
            with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
                text = f.read()
        except Exception:
            # Fallback to latin-1
            try:
                with open(file_path, "rb") as f_bin:
                    text = f_bin.read().decode('latin-1', errors='replace')
            except Exception as e:
                logger.error(f"Text file processing error for {filename}: {e}")
                raise
        
        if text.strip():
            segments.append({
                "text": text,
                "metadata": {
                    "filename": filename,
                    "type": os.path.splitext(filename)[1].lower()
                }
            })
        
        return segments
    
    def split_text_semantic(self, text: str, max_chars: Optional[int] = None) -> List[str]:
        """
        Split text into semantic chunks
        
        Args:
            text: Text to split
            max_chars: Maximum characters per chunk (uses config default if None)
            
        Returns:
            List of text chunks
        """
        max_chars = max_chars or self.max_chunk_chars
        logger.debug(f"Splitting text (length {len(text)}) semantically, max_chars {max_chars}")
        
        def _split_by_delimiters(text_block, delimiters_regex):
            parts = re.split(delimiters_regex, text_block)
            result = []
            for i in range(0, len(parts), 2):
                part_text = parts[i]
                delimiter = parts[i+1] if (i+1 < len(parts) and parts[i+1]) else ''
                if part_text or delimiter:
                    result.append(part_text + delimiter)
            return [p.strip() for p in result if p.strip()]
        
        # Split into initial chunks (paragraphs and sentences)
        initial_chunks = []
        paragraphs = re.split(r'\n\s*\n+', text.strip())
        
        for para in paragraphs:
            if not para.strip():
                continue
            sentences = _split_by_delimiters(para, r'([.!?])\s+')
            initial_chunks.extend(sentences)
        
        # Combine into final chunks respecting max_chars
        final_chunks = []
        current_chunk = ""
        
        for chunk_part in initial_chunks:
            if not chunk_part:
                continue
            
            if len(current_chunk) + len(chunk_part) + 1 <= max_chars:
                current_chunk = chunk_part if not current_chunk else current_chunk + " " + chunk_part
            else:
                if current_chunk:
                    final_chunks.append(current_chunk.strip())
                
                if len(chunk_part) > max_chars:
                    logger.debug(f"Hard splitting oversized chunk (len {len(chunk_part)})")
                    for i in range(0, len(chunk_part), max_chars):
                        final_chunks.append(chunk_part[i:i+max_chars].strip())
                    current_chunk = ""
                else:
                    current_chunk = chunk_part
        
        if current_chunk:
            final_chunks.append(current_chunk.strip())
        
        result = [c for c in final_chunks if c]
        self.stats["total_chunks"] += len(result)
        logger.debug(f"Text split into {len(result)} chunks")
        
        return result
    
    def load_documents_from_folder(
        self,
        folder_path: Optional[str] = None,
        progress_callback: Optional[Callable[[float, str], None]] = None
    ) -> List[Document]:
        """
        Load all supported files from a folder and return a flat list of Documents.
        """
        import time

        folder_path = folder_path or self.files_directory
        logger.info(f"Loading documents from folder: {folder_path}")

        all_docs: List[Document] = []
        total_size_bytes, paths = 0, []

        # Find all supported files
        for root, _, files in os.walk(folder_path):
            for f_name in files:
                ext = os.path.splitext(f_name)[1].lower()
                if ext in self.supported_extensions:
                    p = os.path.join(root, f_name)
                    paths.append(p)
                    try:
                        total_size_bytes += os.path.getsize(p)
                    except OSError as e:
                        logger.warning(f"Could not get size of {p}: {e}")

        if not paths:
            logger.warning("No supported files found in the document directory.")
            if progress_callback:
                progress_callback(1.0, "No files found")
            return []

        logger.info(f"Found {len(paths)} files to process, total size {total_size_bytes/1e6:.2f} MB")

        for i, p in enumerate(paths):
            logger.info(f"Processing file {i+1}/{len(paths)}: {p}")
            if progress_callback:
                progress = (i + 1) / len(paths)
                status = f"Processing {os.path.basename(p)}..."
                progress_callback(progress, status)

            docs = self.process_file(p)
            all_docs.extend(docs)

        logger.info(f"Loaded {len(all_docs)} chunks as Document objects")
        if progress_callback:
            progress_callback(1.0, f"Completed: {len(all_docs)} chunks loaded")
        return all_docs

    
    def process_file(self, file_path: str) -> List[Document]:
        """
        Process a single file and return a list of Document objects.
        Each Document contains one chunk in `page_content` and a rich `metadata`.
        """
        documents: List[Document] = []

        segments = self.extract_text_with_fallback(file_path)
        for seg_idx, seg in enumerate(segments):
            txt = seg["text"].strip()
            if not txt:
                continue

            seg_chunks = self.split_text_semantic(txt)
            for chunk_idx, chunk in enumerate(seg_chunks):
                # Build source identifier and metadata
                meta = dict(seg.get("metadata", {}))  # copy
                # Add contextual metadata for traceability
                meta["segment_index"] = seg_idx
                meta["chunk_index"] = chunk_idx
                # Optional human-readable source label
                src_label = meta.get("filename", "unknown")
                for key in ["page", "slide", "sheet"]:
                    if key in meta:
                        src_label += f" - {key.title()} {meta[key]}"
                if len(segments) > 1:
                    src_label += f" (Segment {seg_idx+1})"
                if len(seg_chunks) > 1:
                    src_label += f" (Chunk {chunk_idx+1})"
                meta["source"] = src_label

                documents.append(Document(page_content=chunk, metadata=meta))

        self.stats["total_chunks"] += len(documents)
        return documents
    
    def get_stats(self) -> Dict[str, Any]:
        """Get processing statistics"""
        return self.stats.copy()
    
    def reset_stats(self):
        """Reset processing statistics"""
        self.stats = {
            "files_processed": 0,
            "extraction_errors": 0,
            "total_chunks": 0,
            "ocr_attempts": 0,
            "ocr_successes": 0
        }